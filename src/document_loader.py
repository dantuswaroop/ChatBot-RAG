import fitz  # PyMuPDF for PDF
from pathlib import Path
import os
from typing import List, Tuple, Dict, Any
import logging

# Document-specific imports
try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    from openpyxl import load_workbook
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    import xlrd
    XLRD_AVAILABLE = True
except ImportError:
    XLRD_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class DocumentLoader:
    """
    A comprehensive document loader that supports multiple file formats:
    - PDF (.pdf)
    - Microsoft Word (.doc, .docx)
    - Microsoft Excel (.xls, .xlsx)
    - Microsoft PowerPoint (.ppt, .pptx)
    """
    
    SUPPORTED_EXTENSIONS = {
        '.pdf': 'extract_pdf',
        '.docx': 'extract_docx',
        '.doc': 'extract_doc',
        '.xlsx': 'extract_xlsx',
        '.xls': 'extract_xls',
        '.pptx': 'extract_pptx',
        '.ppt': 'extract_ppt'
    }
    
    def __init__(self):
        self._check_dependencies()
    
    def _check_dependencies(self):
        """Check which optional dependencies are available"""
        missing_deps = []
        
        if not DOCX_AVAILABLE:
            missing_deps.append("python-docx (for .docx files)")
        if not OPENPYXL_AVAILABLE:
            missing_deps.append("openpyxl (for .xlsx files)")
        if not XLRD_AVAILABLE:
            missing_deps.append("xlrd (for .xls files)")
        if not PPTX_AVAILABLE:
            missing_deps.append("python-pptx (for .pptx files)")
        
        if missing_deps:
            logger.warning(f"Missing optional dependencies: {', '.join(missing_deps)}")
            logger.warning("Some file formats may not be supported. Install missing packages with pip.")
    
    def extract_text_from_documents(self, doc_dir: str) -> List[Tuple[str, str, List[Dict]]]:
        """
        Extract text from documents in the specified directory.
        Returns: list of (filename, full_text, page_metadata)
        """
        results = []
        doc_dir = Path(doc_dir)
        
        if not doc_dir.exists():
            logger.error(f"Directory does not exist: {doc_dir}")
            return results
        
        # Get all supported files recursively
        supported_files = []
        for ext in self.SUPPORTED_EXTENSIONS.keys():
            # Use ** for recursive search through all subdirectories
            supported_files.extend(doc_dir.glob(f"**/*{ext}"))
        
        if not supported_files:
            logger.warning(f"No supported document files found in {doc_dir} (searched recursively)")
            return results
        
        logger.info(f"Found {len(supported_files)} supported document(s) to process")
        
        for file_path in supported_files:
            try:
                result = self._extract_from_file(file_path)
                if result:
                    results.append(result)
            except Exception as e:
                logger.error(f"Error processing {file_path.name}: {str(e)}")
                continue
        
        return results
    
    def _extract_from_file(self, file_path: Path) -> Tuple[str, str, List[Dict]] or None:
        """Extract text from a single file based on its extension"""
        # Check file size
        file_size = os.path.getsize(file_path)
        if file_size == 0:
            logger.warning(f"Skipping empty file: {file_path.name}")
            return None
        elif file_size < 50:  # Very small files are likely corrupted
            logger.warning(f"Skipping potentially corrupted file: {file_path.name} (size: {file_size} bytes)")
            return None
        
        ext = file_path.suffix.lower()
        if ext not in self.SUPPORTED_EXTENSIONS:
            logger.warning(f"Unsupported file type: {ext}")
            return None
        
        extract_method = getattr(self, self.SUPPORTED_EXTENSIONS[ext])
        return extract_method(file_path)
    
    def extract_pdf(self, file_path: Path) -> Tuple[str, str, List[Dict]] or None:
        """Extract text from PDF files using PyMuPDF"""
        try:
            doc = fitz.open(file_path)
            full_text = ""
            page_metadata = []
            
            for page_num, page in enumerate(doc, 1):
                page_text = page.get_text()
                if page_text.strip():  # Only add non-empty pages
                    full_text += page_text + "\n"
                    
                    page_metadata.append({
                        "page_number": page_num,
                        "text": page_text,
                        "char_start": len(full_text) - len(page_text) - 1,
                        "char_end": len(full_text) - 1
                    })
            
            doc.close()
            
            if full_text.strip():
                logger.info(f"✅ Processed PDF: {file_path.name} ({len(full_text)} characters, {len(page_metadata)} pages)")
                return (file_path.stem, full_text, page_metadata)
            else:
                logger.warning(f"⚠️ No extractable text found in PDF: {file_path.name}")
                return None
                
        except Exception as e:
            logger.error(f"❌ Error processing PDF {file_path.name}: {str(e)}")
            return None
    
    def extract_docx(self, file_path: Path) -> Tuple[str, str, List[Dict]] or None:
        """Extract text from DOCX files"""
        if not DOCX_AVAILABLE:
            logger.error("python-docx not available. Cannot process .docx files.")
            return None
        
        try:
            doc = DocxDocument(file_path)
            full_text = ""
            page_metadata = []
            
            for i, paragraph in enumerate(doc.paragraphs):
                para_text = paragraph.text
                if para_text.strip():
                    full_text += para_text + "\n"
                    
                    page_metadata.append({
                        "paragraph_number": i + 1,
                        "text": para_text,
                        "char_start": len(full_text) - len(para_text) - 1,
                        "char_end": len(full_text) - 1
                    })
            
            if full_text.strip():
                logger.info(f"✅ Processed DOCX: {file_path.name} ({len(full_text)} characters, {len(page_metadata)} paragraphs)")
                return (file_path.stem, full_text, page_metadata)
            else:
                logger.warning(f"⚠️ No extractable text found in DOCX: {file_path.name}")
                return None
                
        except Exception as e:
            logger.error(f"❌ Error processing DOCX {file_path.name}: {str(e)}")
            return None
    
    def extract_doc(self, file_path: Path) -> Tuple[str, str, List[Dict]] or None:
        """Extract text from DOC files using PyMuPDF (which can handle DOC files)"""
        try:
            # PyMuPDF can handle some DOC files
            doc = fitz.open(file_path)
            full_text = ""
            page_metadata = []
            
            for page_num, page in enumerate(doc, 1):
                page_text = page.get_text()
                if page_text.strip():
                    full_text += page_text + "\n"
                    
                    page_metadata.append({
                        "page_number": page_num,
                        "text": page_text,
                        "char_start": len(full_text) - len(page_text) - 1,
                        "char_end": len(full_text) - 1
                    })
            
            doc.close()
            
            if full_text.strip():
                logger.info(f"✅ Processed DOC: {file_path.name} ({len(full_text)} characters, {len(page_metadata)} pages)")
                return (file_path.stem, full_text, page_metadata)
            else:
                logger.warning(f"⚠️ No extractable text found in DOC: {file_path.name}")
                return None
                
        except Exception as e:
            logger.error(f"❌ Error processing DOC {file_path.name}: {str(e)}")
            logger.info("Note: For better DOC support, consider converting to DOCX format")
            return None
    
    def extract_xlsx(self, file_path: Path) -> Tuple[str, str, List[Dict]] or None:
        """Extract text from XLSX files"""
        if not OPENPYXL_AVAILABLE:
            logger.error("openpyxl not available. Cannot process .xlsx files.")
            return None
        
        try:
            workbook = load_workbook(file_path, data_only=True)
            full_text = ""
            page_metadata = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_text = f"=== Sheet: {sheet_name} ===\n"
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        sheet_text += row_text + "\n"
                
                if sheet_text.strip():
                    full_text += sheet_text + "\n"
                    
                    page_metadata.append({
                        "sheet_name": sheet_name,
                        "text": sheet_text,
                        "char_start": len(full_text) - len(sheet_text) - 1,
                        "char_end": len(full_text) - 1
                    })
            
            workbook.close()
            
            if full_text.strip():
                logger.info(f"✅ Processed XLSX: {file_path.name} ({len(full_text)} characters, {len(page_metadata)} sheets)")
                return (file_path.stem, full_text, page_metadata)
            else:
                logger.warning(f"⚠️ No extractable text found in XLSX: {file_path.name}")
                return None
                
        except Exception as e:
            logger.error(f"❌ Error processing XLSX {file_path.name}: {str(e)}")
            return None
    
    def extract_xls(self, file_path: Path) -> Tuple[str, str, List[Dict]] or None:
        """Extract text from XLS files"""
        if not XLRD_AVAILABLE:
            logger.error("xlrd not available. Cannot process .xls files.")
            return None
        
        try:
            import xlrd
            workbook = xlrd.open_workbook(file_path)
            full_text = ""
            page_metadata = []
            
            for sheet_idx in range(workbook.nsheets):
                sheet = workbook.sheet_by_index(sheet_idx)
                sheet_name = sheet.name
                sheet_text = f"=== Sheet: {sheet_name} ===\n"
                
                for row_idx in range(sheet.nrows):
                    row_data = []
                    for col_idx in range(sheet.ncols):
                        cell = sheet.cell(row_idx, col_idx)
                        row_data.append(str(cell.value) if cell.value else "")
                    
                    row_text = "\t".join(row_data)
                    if row_text.strip():
                        sheet_text += row_text + "\n"
                
                if sheet_text.strip():
                    full_text += sheet_text + "\n"
                    
                    page_metadata.append({
                        "sheet_name": sheet_name,
                        "text": sheet_text,
                        "char_start": len(full_text) - len(sheet_text) - 1,
                        "char_end": len(full_text) - 1
                    })
            
            if full_text.strip():
                logger.info(f"✅ Processed XLS: {file_path.name} ({len(full_text)} characters, {len(page_metadata)} sheets)")
                return (file_path.stem, full_text, page_metadata)
            else:
                logger.warning(f"⚠️ No extractable text found in XLS: {file_path.name}")
                return None
                
        except Exception as e:
            logger.error(f"❌ Error processing XLS {file_path.name}: {str(e)}")
            return None
    
    def extract_pptx(self, file_path: Path) -> Tuple[str, str, List[Dict]] or None:
        """Extract text from PPTX files"""
        if not PPTX_AVAILABLE:
            logger.error("python-pptx not available. Cannot process .pptx files.")
            return None
        
        try:
            presentation = Presentation(file_path)
            full_text = ""
            page_metadata = []
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_text = f"=== Slide {slide_num} ===\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                
                if slide_text.strip():
                    full_text += slide_text + "\n"
                    
                    page_metadata.append({
                        "slide_number": slide_num,
                        "text": slide_text,
                        "char_start": len(full_text) - len(slide_text) - 1,
                        "char_end": len(full_text) - 1
                    })
            
            if full_text.strip():
                logger.info(f"✅ Processed PPTX: {file_path.name} ({len(full_text)} characters, {len(page_metadata)} slides)")
                return (file_path.stem, full_text, page_metadata)
            else:
                logger.warning(f"⚠️ No extractable text found in PPTX: {file_path.name}")
                return None
                
        except Exception as e:
            logger.error(f"❌ Error processing PPTX {file_path.name}: {str(e)}")
            return None
    
    def extract_ppt(self, file_path: Path) -> Tuple[str, str, List[Dict]] or None:
        """Extract text from PPT files using PyMuPDF"""
        try:
            # PyMuPDF can handle some PPT files
            doc = fitz.open(file_path)
            full_text = ""
            page_metadata = []
            
            for page_num, page in enumerate(doc, 1):
                page_text = page.get_text()
                if page_text.strip():
                    slide_text = f"=== Slide {page_num} ===\n" + page_text + "\n"
                    full_text += slide_text
                    
                    page_metadata.append({
                        "slide_number": page_num,
                        "text": page_text,
                        "char_start": len(full_text) - len(slide_text),
                        "char_end": len(full_text)
                    })
            
            doc.close()
            
            if full_text.strip():
                logger.info(f"✅ Processed PPT: {file_path.name} ({len(full_text)} characters, {len(page_metadata)} slides)")
                return (file_path.stem, full_text, page_metadata)
            else:
                logger.warning(f"⚠️ No extractable text found in PPT: {file_path.name}")
                return None
                
        except Exception as e:
            logger.error(f"❌ Error processing PPT {file_path.name}: {str(e)}")
            logger.info("Note: For better PPT support, consider converting to PPTX format")
            return None


# Backward compatibility function
def extract_text_from_pdfs(pdf_dir: str) -> List[Tuple[str, str, List[Dict]]]:
    """
    Backward compatibility function for existing code.
    Now extracts from all supported document formats, not just PDFs.
    """
    loader = DocumentLoader()
    return loader.extract_text_from_documents(pdf_dir)


# Main function for testing
if __name__ == "__main__":
    import sys
    
    if len(sys.argv) != 2:
        print("Usage: python document_loader.py <directory_path>")
        sys.exit(1)
    
    directory = sys.argv[1]
    loader = DocumentLoader()
    results = loader.extract_text_from_documents(directory)
    
    print(f"\nProcessed {len(results)} documents:")
    for filename, text, metadata in results:
        print(f"- {filename}: {len(text)} characters, {len(metadata)} sections")